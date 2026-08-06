import io

from pypdf import PdfReader

from docx import Document

from pptx import Presentation


def extract_text(filename: str, content: bytes) -> str:

    ext = filename.lower().rsplit(".", 1)[-1]

    if ext == "pdf":

        return _extract_pdf(content)

    elif ext == "docx":

        return _extract_docx(content)

    elif ext == "pptx":

        return _extract_pptx(content)

    elif ext == "txt":

        return content.decode("utf-8", errors="ignore")

    else:

        raise ValueError(f"Unsupported file type: {ext}")


def _extract_pdf(content: bytes) -> str:

    reader = PdfReader(io.BytesIO(content))

    return "\n".join(

        page.extract_text() or ""

        for page in reader.pages

    )


def _extract_docx(content: bytes) -> str:

    doc = Document(io.BytesIO(content))

    return "\n".join(

        p.text

        for p in doc.paragraphs

    )


def _extract_pptx(content: bytes) -> str:

    prs = Presentation(io.BytesIO(content))

    parts = []

    for slide in prs.slides:

        for shape in slide.shapes:

            if shape.has_text_frame:

                parts.append(shape.text_frame.text)

    return "\n".join(parts)